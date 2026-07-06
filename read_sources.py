from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


ROOT = Path(__file__).resolve().parent.parent


def read_docx(path: Path) -> None:
    doc = Document(path)
    print(f"\n===== WORD: {path.name} =====")
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            print(text)
    for table_index, table in enumerate(doc.tables, start=1):
        print(f"--- Table {table_index} ---")
        for row in table.rows:
            print(" | ".join(cell.text.replace("\n", " / ").strip() for cell in row.cells))


def read_xlsx(path: Path) -> None:
    workbook = load_workbook(path, data_only=False, read_only=True)
    print(f"\n===== EXCEL: {path.name} =====")
    for sheet in workbook.worksheets:
        print(f"--- Sheet: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(values):
                print(" | ".join(values))


def read_pptx(path: Path) -> None:
    presentation = Presentation(path)
    print(f"\n===== POWERPOINT: {path.name} =====")
    for slide_index, slide in enumerate(presentation.slides, start=1):
        print(f"--- Slide {slide_index} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(shape.text.strip().replace("\n", " | "))
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    print(" | ".join(cell.text.strip() for cell in row.cells))


for source in sorted(ROOT.iterdir()):
    suffix = source.suffix.lower()
    if suffix == ".docx":
        read_docx(source)
    elif suffix == ".xlsx":
        read_xlsx(source)
    elif suffix == ".pptx":
        read_pptx(source)
